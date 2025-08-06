import os
import csv
import fitz  # PyMuPDF for PDF files
import tiktoken
import tkinter as tk
from tkinter import filedialog, messagebox
from docx import Document  # For .docx files
from openpyxl import load_workbook  # For .xlsx files
from pptx import Presentation  # For .pptx files
import threading

# Initialize a variable to control the dot effect
dot_count = 0
processing = False

def extract_text_from_pdf(pdf_path):
    """Extracts text from each page of a PDF."""
    text = ""
    with fitz.open(pdf_path) as doc:
        for page in doc:
            text += page.get_text("text")
    return text

def extract_text_from_docx(docx_path):
    """Extracts text from a .docx (Word) file."""
    doc = Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_xlsx(xlsx_path):
    """Extracts text from a .xlsx (Excel) file."""
    workbook = load_workbook(xlsx_path, data_only=True)
    text = ""
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            row_text = " ".join([str(cell.value) if cell.value is not None else "" for cell in row])
            text += row_text + "\n"
    return text

def extract_text_from_pptx(pptx_path):
    """Extracts text from a .pptx (PowerPoint) file."""
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def count_tokens(text, model="gpt-4"):
    """Counts the number of tokens in the given text using the specified model's tokenizer."""
    encoding = tiktoken.encoding_for_model(model)
    tokens = encoding.encode(text)
    return len(tokens)


def update_processing_dots():
    """Updates the processing label with dynamic dots."""
    global dot_count
    if processing:
        # Cycle through 1, 2, and 3 dots
        dot_count = (dot_count + 1) % 4
        dots = '.' * dot_count
        result_label.config(text=f"Processing{dots}")
        root.after(500, update_processing_dots)  # Update every 500 milliseconds

def process_file(file_path):
    """Process the selected file in a separate thread."""
    global processing

    try:
        # Determine the file type and extract text accordingly
        if file_path.endswith(".pdf"):
            extracted_text = extract_text_from_pdf(file_path)
        elif file_path.endswith(".docx"):
            extracted_text = extract_text_from_docx(file_path)
        elif file_path.endswith(".xlsx"):
            extracted_text = extract_text_from_xlsx(file_path)
        elif file_path.endswith(".pptx"):
            extracted_text = extract_text_from_pptx(file_path)
        else:
            messagebox.showerror("Unsupported file", "The selected file type is not supported.")
            result_label.config(text="Total tokens:")
            processing = False
            return

        # Calculate and display token count for the extracted text
        token_count = count_tokens(extracted_text)
        result_entry.delete(0, tk.END)
        result_entry.insert(0, str(token_count))
        result_label.config(text="Total tokens:")
    except Exception as e:
        result_label.config(text="An error occurred:")
        result_entry.delete(0, tk.END)
        result_entry.insert(0, str(e))
    finally:
        processing = False  # Stop the dot effect

def open_file_and_analyze():
    global processing

    # Open file dialog to select a file
    file_path = filedialog.askopenfilename(filetypes=[("Supported files", "*.pdf *.docx *.xlsx *.pptx")])
    if not file_path:
        return

    # Display the truncated file name in the file_label (limit to 40 characters)
    file_name = file_path.split("/")[-1]
    if len(file_name) > 40:
        file_name = file_name[:37] + "..."
    file_label.config(text=f"Selected File: {file_name}")

    # Start the processing animation and file processing
    processing = True
    update_processing_dots()  # Start the dot effect
    threading.Thread(target=process_file, args=(file_path,), daemon=True).start()




def process_folder(folder_path):
    global processing
    results = []
    total_tokens = 0
    supported_ext = (".pdf", ".docx", ".xlsx", ".pptx")
    try:
        for fname in os.listdir(folder_path):
            if fname.lower().endswith(supported_ext):
                fpath = os.path.join(folder_path, fname)
                if fname.endswith(".pdf"):
                    text = extract_text_from_pdf(fpath)
                elif fname.endswith(".docx"):
                    text = extract_text_from_docx(fpath)
                elif fname.endswith(".xlsx"):
                    text = extract_text_from_xlsx(fpath)
                elif fname.endswith(".pptx"):
                    text = extract_text_from_pptx(fpath)
                else:
                    continue
                tokens = count_tokens(text)
                results.append([fname, tokens])
                total_tokens += tokens
        # Save results to CSV only if checkbox is checked
        if create_csv_var.get():
            csv_path = os.path.join(folder_path, "token_counts.csv")
            with open(csv_path, "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f)
                writer.writerow(["filename", "tokencount"])
                writer.writerows(results)
        result_entry.delete(0, tk.END)
        result_entry.insert(0, str(total_tokens))
        result_label.config(text="Total tokens in folder:")
    except Exception as e:
        result_label.config(text="An error occurred:")
        result_entry.delete(0, tk.END)
        result_entry.insert(0, str(e))
    finally:
        processing = False





def open_folder_and_analyze():
    global processing
    folder_path = filedialog.askdirectory()
    if not folder_path:
        return
    folder_name = os.path.basename(folder_path)
    file_label.config(text=f"Selected Folder: {folder_name}")
    processing = True
    update_processing_dots()
    threading.Thread(target=process_folder, args=(folder_path,), daemon=True).start()


# Set up the Tkinter interface
root = tk.Tk()
root.title("File Token Counter")
root.geometry("")  # Let the window size dynamically adjust
root.resizable(False, False)  # Lock the window size
root.minsize(370, 250)  # Set minimum window size



# Label to display the selected folder name with adjustable wraplength
file_label = tk.Label(root, text="Selected Folder: None", font=("Arial", 10), wraplength=400, anchor="w", justify="left")
file_label.pack(pady=5, fill="x", padx=10)

# Frame for the token count result
frame = tk.Frame(root, bd=2, relief="sunken", padx=5, pady=5)
frame.pack(pady=10, fill="x", padx=10)

# Label for token count description
result_label = tk.Label(frame, text="Total tokens:", font=("Arial", 10), anchor="w", justify="left")
result_label.pack(side="left")

# Entry widget to display and copy only the token count
result_entry = tk.Entry(frame, relief="flat", font=("Arial", 10), state="normal", readonlybackground="white", width=10)
result_entry.pack(side="right", fill="x")
create_csv_var = tk.BooleanVar(value=True)
csv_checkbox = tk.Checkbutton(root, text="Create CSV file", variable=create_csv_var)
csv_checkbox.pack(pady=2)
# Button to select file and count tokens
select_button = tk.Button(root, text="Select Folder and Analyze", command=open_folder_and_analyze)
select_button.pack(pady=10)
# Exit button to close the application
exit_button = tk.Button(root, text="Exit", command=root.quit)
exit_button.pack(pady=10)

# Version number label in the bottom right
version_label = tk.Label(root, text="Version 1.1", font=("Arial", 8), anchor="se")
version_label.place(relx=1.0, rely=1.0, x=-5, y=-5, anchor="se")


# Run the Tkinter main loop
root.mainloop()
